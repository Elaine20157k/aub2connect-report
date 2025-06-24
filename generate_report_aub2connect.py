UPLOAD_FOLDER = "uploads"

import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.patches import Circle
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def generate_ppt_report(excel_path, output_path, logo_path):
    df = pd.read_excel(excel_path)

    total = len(df)
    checked_in = df['check-in'].str.upper().value_counts().get('Y', 0)
    not_checked_in = df['check-in'].str.upper().value_counts().get('N', 0)
    attendance_rate = round((checked_in / total) * 100, 2)

    job_counts = df['title'].value_counts()
    company_counts = df['company'].value_counts()
    no_show_df = pd.DataFrame()
    if 'check-in' in df.columns:
        no_show_df = df[df['check-in'].str.upper() == 'N']

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_slide_layout)

    # 添加 Logo
    left = Inches(0.2)
    top = Inches(0.2)
    height = Inches(0.8)
    slide.shapes.add_picture(logo_path, left, top, height=height)

    # 添加文字框
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Event Dashboard Summary"
    p.font.size = Pt(28)
    p.font.bold = True

    # 饼图
    fig, ax = plt.subplots(figsize=(2, 2))
    sizes = [checked_in, not_checked_in]
    labels = ['Checked In', 'Not Checked In']
    ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90)
    ax.axis('equal')
    plt.tight_layout()
    chart_path = os.path.join("charts", "pie.png")
    if not os.path.exists("charts"):
        os.makedirs("charts")
    plt.savefig(chart_path)
    plt.close()

    slide.shapes.add_picture(chart_path, Inches(1), Inches(2), height=Inches(2))

    prs.save(output_path)
