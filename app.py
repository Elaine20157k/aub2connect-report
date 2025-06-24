import os
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Circle
from pptx import Presentation
from pptx.util import Inches, Pt
from flask import Flask, request, render_template_string, send_file

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
CHART_FOLDER = "charts"

# 确保目录存在
for folder in [UPLOAD_FOLDER, OUTPUT_FOLDER, CHART_FOLDER]:
    if os.path.exists(folder) and not os.path.isdir(folder):
        os.remove(folder)
    os.makedirs(folder, exist_ok=True)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

HTML_FORM = """
<!doctype html>
<title>Upload Excel</title>
<h1>Upload Excel File to Generate PPT Report</h1>
<form method=post enctype=multipart/form-data action="/upload">
  <input type=file name=file required>
  <input type=submit value="Upload and Generate">
</form>
"""

@app.route("/", methods=["GET"])
def index():
    return render_template_string(HTML_FORM)

@app.route("/upload", methods=["POST"])
def upload_file():
    file = request.files.get("file")
    if not file or file.filename == '':
        return "No file selected.", 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(OUTPUT_FOLDER, "report.pptx")

    file.save(input_path)

    try:
        generate_report(input_path, output_path)
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return f"Error generating report: {str(e)}", 500
def generate_report(excel_path, output_path):
    df = pd.read_excel(excel_path)

    total = len(df)
    checked_in = df["check-in"].str.upper().value_counts().get("Y", 0)
    not_checked_in = df["check-in"].str.upper().value_counts().get("N", 0)
    attendance_rate = round((checked_in / total) * 100, 2)

    job_counts = df["title"].value_counts()
    company_counts = df["company"].value_counts()
    no_show_df = df[df["check-in"].str.upper() == "N"] if "check-in" in df.columns else pd.DataFrame()

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # KPI boxes
    kpi_values = [
        ("Total Registered", total),
        ("Checked In", checked_in),
        ("Attendance Rate", f"{attendance_rate}%"),
        ("No-shows", len(no_show_df))
    ]
    for i, (label, value) in enumerate(kpi_values):
        left = Inches(0.5 + i * 3.0)
        top = Inches(0.3)
        width = Inches(2.7)
        height = Inches(0.8)
        box = slide.shapes.add_shape(
            autoshape_type_id=1, left=left, top=top, width=width, height=height
        )
        tf = box.text_frame
        tf.text = f"{label}\n{value}"
        for p in tf.paragraphs:
            p.font.size = Pt(14)
        box.fill.solid()
        box.fill.fore_color.rgb = (0, 102, 204)

    # Pie chart
    pie_path = os.path.join(CHART_FOLDER, "pie.png")
    fig, ax = plt.subplots()
    ax.pie([checked_in, not_checked_in],
           labels=["Checked In", "Not Checked"],
           autopct="%1.1f%%",
           startangle=90)
    centre_circle = Circle((0, 0), 0.70, fc="white")
    fig.gca().add_artist(centre_circle)
    plt.axis('equal')
    plt.tight_layout()
    plt.savefig(pie_path)
    plt.close()
    slide.shapes.add_picture(pie_path, Inches(0.5), Inches(1.3), height=Inches(2.5))

    # Job title bar chart
    bar_path = os.path.join(CHART_FOLDER, "bar.png")
    fig2, ax2 = plt.subplots()
    job_counts.head(6).plot(kind='bar', ax=ax2)
    plt.tight_layout()
    plt.savefig(bar_path)
    plt.close()
    slide.shapes.add_picture(bar_path, Inches(3.8), Inches(1.3), height=Inches(2.5))

    # No-show table
    if not no_show_df.empty:
        table_data = no_show_df[["company", "name", "email", "phone"]].values.tolist()
        rows, cols = len(table_data) + 1, 4
        left = Inches(6.7)
        top = Inches(1.3)
        width = Inches(6)
        height = Inches(2.5)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        headers = ["Company", "Name", "Email", "Phone"]
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
        for i, row in enumerate(table_data):
            for j, val in enumerate(row):
                table.cell(i + 1, j).text = str(val)

    # Summary
    summary_text = f"A total of {total} participants registered, with {checked_in} checked in ({attendance_rate}%)."
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12), Inches(1.2))
    tf = txBox.text_frame
    tf.text = "Summary"
    tf.add_paragraph().text = summary_text

    # Footer
    footer = slide.shapes.add_textbox(Inches(11), Inches(6.9), Inches(2), Inches(0.3))
    footer.text_frame.text = "AuB2Connect"

    prs.save(output_path)
